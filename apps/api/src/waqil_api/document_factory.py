"""The document factory — turns an authored outline into a real file.

No model emits a .pptx or .pdf; binary formats are not token streams. So the
work is split, and the split is the whole design:

* The **model** writes a ``DocumentOutlineV1`` — title, sections, bullets, and
  markdown tables. Content only, which every provider in the lineup can do.
* The **host** renders it with the fixed, tested code below. Nothing here is
  model-authored, so nothing here needs the sandbox that contains authored
  code; the risk a sandbox exists to hold simply is not present.

That is why a deck works the same on Command A+, Grok, gpt-oss, or a small
local model: the fragile half was never given to the model. The alternative —
having the model write python-pptx code per request — is exactly the shape the
build evals showed models fumbling.

Rendering degrades rather than fails: a table with no columns is skipped, an
empty section becomes a heading, and a document with no sections at all still
produces a valid title page.
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from io import BytesIO

from .contracts import DocumentOutlineV1

# The house style. One palette, one type scale, applied by the renderer so no
# two documents can disagree about what Metis output looks like.
INK = (0x21, 0x1F, 0x1D)
MUTED = (0x6B, 0x66, 0x60)
ACCENT = (0x7A, 0x54, 0xA8)
PAPER = (0xFA, 0xF8, 0xF4)

PPTX_MEDIA_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
PDF_MEDIA_TYPE = "application/pdf"

# Which format an explicit request is asking for. Only unambiguous FILE words
# count. "brief", "report", "write-up", and "summary" describe content people
# want in the chat far more often than they describe a file — reading them as
# file requests turns "give me a short brief on these models" into a surprise
# download, which is the opposite of helpful.
_PPTX_WORDS = re.compile(
    r"\b(?:deck|slides?|slide\s*deck|presentation|powerpoint|pptx?|keynote)\b",
    re.IGNORECASE,
)
_PDF_WORDS = re.compile(
    r"\b(?:pdf|one[-\s]?pagers?|onepagers?|hand[-\s]?outs?)\b", re.IGNORECASE
)
_FORMAT_NOUN = (
    r"(?:deck|slides?|presentation|powerpoint|pptx?|keynote|pdf|"
    r"one[-\s]?pagers?|onepagers?|hand[-\s]?outs?)"
)
# The ask itself, in two shapes: a verb of creation aimed at a format noun
# ("build me a deck"), or the export framing ("summarize this as a pdf").
# Both require the format to be named, so a question *about* a deck never
# renders one.
_DOCUMENT_REQUEST = re.compile(
    r"\b(?:make|build|create|generate|draft|write|put\s+together|prepare|"
    r"turn\s+(?:this|that|it)\s+into|export|produce|give\s+me|send\s+me)\b"
    r"[^.?!\n]{0,60}\b" + _FORMAT_NOUN + r"\b"
    r"|\bas\s+an?\s+" + _FORMAT_NOUN + r"\b",
    re.IGNORECASE,
)


class DocumentRenderError(RuntimeError):
    """The outline could not be rendered into a file."""


def is_explicit_document_request(prompt: str) -> bool:
    """True when the prompt asks for a document to be produced.

    Conservative on purpose, in the same spirit as the toolify and web
    signals: the user asking a question about decks must never cause one to
    be built."""
    return bool(_DOCUMENT_REQUEST.search(prompt))


def requested_format(prompt: str) -> str:
    """"pptx" or "pdf" for this prompt. Slides win ties — a request naming
    both ("a deck I can send as a PDF") is still fundamentally a deck."""
    if _PPTX_WORDS.search(prompt):
        return "pptx"
    if _PDF_WORDS.search(prompt):
        return "pdf"
    return "pdf"


def filename_for(title: str, document_format: str) -> str:
    """A safe, readable filename derived from the document's own title."""
    slug = re.sub(r"[^A-Za-z0-9]+", "-", title or "metis-document").strip("-").lower()
    return f"{slug[:60] or 'metis-document'}.{document_format}"


# A markdown table: pipe-delimited rows, with a dashed divider under the
# header. Written by the model inside `body`, parsed back out here.
_TABLE_LINE = re.compile(r"^\s*\|.*\|\s*$")
_TABLE_DIVIDER = re.compile(r"^\s*\|[\s:|-]+\|\s*$")


def _split_cells(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


@dataclass(frozen=True, slots=True)
class TextBlock:
    text: str


@dataclass(frozen=True, slots=True)
class TableBlock:
    columns: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)


BodyBlock = TextBlock | TableBlock


def split_body(body: str) -> list[BodyBlock]:
    """A section's body as ordered prose and table blocks.

    Ragged rows are padded to the header width rather than dropped: a table
    that silently loses a column is worse than one with a blank cell."""
    blocks: list[BodyBlock] = []
    prose: list[str] = []
    table: list[list[str]] = []

    def flush_prose() -> None:
        text = "\n".join(prose).strip()
        if text:
            blocks.append(TextBlock(text))
        prose.clear()

    def flush_table() -> None:
        if table:
            columns, *rows = table
            width = len(columns)
            blocks.append(
                TableBlock(columns, [(row + [""] * width)[:width] for row in rows])
            )
        table.clear()

    for line in body.splitlines():
        if _TABLE_LINE.match(line):
            if _TABLE_DIVIDER.match(line):
                continue  # the header rule carries no content
            flush_prose()
            table.append(_split_cells(line))
            continue
        flush_table()
        prose.append(line)
    flush_prose()
    flush_table()
    return blocks


def render(outline: DocumentOutlineV1, document_format: str) -> bytes:
    """Render an outline to file bytes. The one entry point."""
    if document_format == "pptx":
        return _render_pptx(outline)
    if document_format == "pdf":
        return _render_pdf(outline)
    raise DocumentRenderError(f"unsupported document format: {document_format}")


def _render_pptx(outline: DocumentOutlineV1) -> bytes:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Emu, Inches, Pt
    except ImportError as exc:  # pragma: no cover - dependency is pinned
        raise DocumentRenderError(
            "python-pptx is required to render slide decks"
        ) from exc

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)  # 16:9, the only sane default
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    def paint_background(slide) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*PAPER)

    def textbox(slide, left, top, width, height):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        return frame

    def style(run, *, size: int, color=INK, bold: bool = False) -> None:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)

    # Title slide — the accent rule under the title is the whole brand here.
    cover = presentation.slides.add_slide(blank)
    paint_background(cover)
    frame = textbox(cover, 0.9, 2.4, 11.5, 1.6)
    style(frame.paragraphs[0].add_run(), size=40, bold=True)
    frame.paragraphs[0].runs[0].text = outline.title or "Untitled"
    if outline.subtitle:
        subtitle = textbox(cover, 0.9, 4.0, 11.5, 1.0)
        run = subtitle.paragraphs[0].add_run()
        run.text = outline.subtitle
        style(run, size=18, color=MUTED)
    rule = cover.shapes.add_shape(1, Inches(0.9), Inches(3.95), Inches(1.6), Emu(38100))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor(*ACCENT)
    rule.line.fill.background()
    rule.shadow.inherit = False

    for section in outline.sections:
        slide = presentation.slides.add_slide(blank)
        paint_background(slide)
        heading = textbox(slide, 0.9, 0.7, 11.5, 1.0)
        run = heading.paragraphs[0].add_run()
        run.text = section.heading or "—"
        style(run, size=28, bold=True)

        top = 1.9
        prose = "\n".join(
            block.text for block in split_body(section.body)
            if isinstance(block, TextBlock)
        )
        if prose:
            body = textbox(slide, 0.9, top, 11.5, 1.2)
            run = body.paragraphs[0].add_run()
            run.text = prose
            style(run, size=15, color=MUTED)
            top += 1.1

        if section.bullets:
            bullets = textbox(slide, 0.9, top, 11.5, 4.0)
            for index, bullet in enumerate(section.bullets):
                paragraph = bullets.paragraphs[0] if index == 0 else bullets.add_paragraph()
                paragraph.space_after = Pt(10)
                run = paragraph.add_run()
                run.text = f"•  {bullet}"
                style(run, size=16)
            top += 0.42 * len(section.bullets) + 0.3

        for block in split_body(section.body):
            if not isinstance(block, TableBlock) or not block.columns or top > 6.6:
                continue
            headers, cells = block.columns, block.rows
            rows, columns = len(cells) + 1, len(headers)
            height = Inches(max(min(0.4 * (rows + 1), 7.5 - top - 0.4), 0.4))
            shape = slide.shapes.add_table(
                rows, columns, Inches(0.9), Inches(top), Inches(11.5), height
            )
            table = shape.table
            for column, label in enumerate(headers):
                cell = table.cell(0, column)
                cell.text = label
                style(cell.text_frame.paragraphs[0].runs[0], size=12, bold=True)
            for row_index, row in enumerate(cells, start=1):
                for column, cell_value in enumerate(row):
                    cell = table.cell(row_index, column)
                    cell.text = cell_value
                    style(cell.text_frame.paragraphs[0].runs[0], size=11)
            top += 0.4 * (rows + 1) + 0.2

        if section.notes:
            slide.notes_slide.notes_text_frame.text = section.notes

    if outline.sources:
        slide = presentation.slides.add_slide(blank)
        paint_background(slide)
        heading = textbox(slide, 0.9, 0.7, 11.5, 0.9)
        run = heading.paragraphs[0].add_run()
        run.text = "Sources"
        style(run, size=28, bold=True)
        frame = textbox(slide, 0.9, 1.8, 11.5, 5.0)
        for index, source in enumerate(outline.sources):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.space_after = Pt(8)
            run = paragraph.add_run()
            run.text = f"[{index + 1}]  {source}"
            style(run, size=12, color=MUTED)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _render_pdf(outline: DocumentOutlineV1) -> bytes:
    try:
        # reportlab ships no py.typed marker; the imports are pinned and the
        # values are used through this module's own typed helpers.
        from reportlab.lib import colors  # type: ignore[import-untyped]
        from reportlab.lib.enums import TA_LEFT  # type: ignore[import-untyped]
        from reportlab.lib.pagesizes import A4  # type: ignore[import-untyped]
        from reportlab.lib.styles import ParagraphStyle  # type: ignore[import-untyped]
        from reportlab.lib.units import mm  # type: ignore[import-untyped]
        from reportlab.platypus import (  # type: ignore[import-untyped]
            HRFlowable,
            ListFlowable,
            ListItem,
            Paragraph,
            SimpleDocTemplate,
            Spacer,
            Table,
            TableStyle,
        )
    except ImportError as exc:  # pragma: no cover - dependency is pinned
        raise DocumentRenderError("reportlab is required to render PDFs") from exc

    def rgb(value: tuple[int, int, int]):
        return colors.Color(value[0] / 255, value[1] / 255, value[2] / 255)

    title_style = ParagraphStyle(
        "MetisTitle", fontName="Helvetica-Bold", fontSize=24, leading=29,
        textColor=rgb(INK), spaceAfter=4,
    )
    subtitle_style = ParagraphStyle(
        "MetisSubtitle", fontName="Helvetica", fontSize=11.5, leading=16,
        textColor=rgb(MUTED), spaceAfter=10,
    )
    heading_style = ParagraphStyle(
        "MetisHeading", fontName="Helvetica-Bold", fontSize=13.5, leading=18,
        textColor=rgb(INK), spaceBefore=16, spaceAfter=6,
    )
    body_style = ParagraphStyle(
        "MetisBody", fontName="Helvetica", fontSize=10.5, leading=15.5,
        textColor=rgb(INK), alignment=TA_LEFT, spaceAfter=6,
    )
    source_style = ParagraphStyle(
        "MetisSource", fontName="Helvetica", fontSize=8.5, leading=12.5,
        textColor=rgb(MUTED),
    )

    def escape(text: str) -> str:
        """Platypus reads a paragraph as mini-HTML, so unescaped user text can
        break the render (or silently vanish) on a stray angle bracket."""
        return (
            text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        )

    story: list = [
        Paragraph(escape(outline.title or "Untitled"), title_style),
        HRFlowable(width="24%", thickness=2.2, color=rgb(ACCENT), spaceAfter=8),
    ]
    if outline.subtitle:
        story.append(Paragraph(escape(outline.subtitle), subtitle_style))

    for section in outline.sections:
        if section.heading:
            story.append(Paragraph(escape(section.heading), heading_style))
        for block in split_body(section.body):
            if isinstance(block, TextBlock):
                story.append(Paragraph(escape(block.text), body_style))
                continue
            if not block.columns:
                continue
            data = [[Paragraph(f"<b>{escape(label)}</b>", body_style) for label in block.columns]]
            data += [[Paragraph(escape(cell), body_style) for cell in row] for row in block.rows]
            table = Table(data, hAlign="LEFT", repeatRows=1)
            table.setStyle(
                TableStyle(
                    [
                        ("LINEBELOW", (0, 0), (-1, 0), 0.9, rgb(ACCENT)),
                        ("LINEBELOW", (0, 1), (-1, -2), 0.4, colors.Color(0.85, 0.84, 0.82)),
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                        ("LEFTPADDING", (0, 0), (-1, -1), 0),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 10),
                        ("TOPPADDING", (0, 0), (-1, -1), 5),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                    ]
                )
            )
            story.extend([table, Spacer(1, 8)])
        if section.bullets:
            story.append(
                ListFlowable(
                    [
                        ListItem(
                            Paragraph(escape(bullet), body_style), leftIndent=12
                        )
                        for bullet in section.bullets
                    ],
                    bulletType="bullet",
                    bulletFontSize=7,
                    bulletColor=rgb(ACCENT),
                    leftIndent=12,
                )
            )
            story.append(Spacer(1, 5))
    if outline.sources:
        story.append(Paragraph("Sources", heading_style))
        for index, source in enumerate(outline.sources, start=1):
            story.append(Paragraph(f"[{index}] {escape(source)}", source_style))

    buffer = BytesIO()
    SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=22 * mm,
        rightMargin=22 * mm,
        topMargin=20 * mm,
        bottomMargin=18 * mm,
        title=outline.title or "Metis document",
        author="Metis",
    ).build(story)
    return buffer.getvalue()
